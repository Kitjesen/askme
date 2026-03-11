"""Add speed comparison slide to Askme v4 Report PPT."""

import sys
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

prs = Presentation("docs/NOVA_Dog_Askme_v4_Report_v5.pptx")

# Add new slide using Blank layout
blank_layout = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(blank_layout)

# Move the new slide (appended at end) to position after slide 13
sldIdLst = prs.slides._sldIdLst
sldId_elements = list(sldIdLst)
last = sldId_elements[-1]
sldIdLst.remove(last)
sldIdLst.insert(13, last)

slide = new_slide
W = prs.slide_width
H = prs.slide_height

# ── Helpers ──────────────────────────────────────────────────

def add_text_box(sl, left, top, width, height, text,
                 font_size=14, bold=False, color=RGBColor(0x33, 0x33, 0x33),
                 alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = sl.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    p.alignment = alignment
    return txBox


def style_cell(cell, text, font_size=11, bold=False,
               color=RGBColor(0x33, 0x33, 0x33), bg_color=None,
               alignment=PP_ALIGN.CENTER, font_name=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    p.alignment = alignment
    cell.text_frame.word_wrap = True
    if bg_color:
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, f"{{{ns}}}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{{{ns}}}srgbClr")
        srgbClr.set("val", bg_color)


# ── Layout ───────────────────────────────────────────────────

margin_left = Inches(0.6)
margin_top = Inches(0.3)
content_width = W - Inches(1.2)

# ── Title ────────────────────────────────────────────────────

add_text_box(slide, margin_left, margin_top, content_width, Inches(0.5),
             "LLM + TTS \u901f\u5ea6\u5bf9\u6bd4 \u2014 MiniMax \u5168\u94fe\u8def\u4f18\u5316",
             font_size=24, bold=True, color=RGBColor(0x1A, 0x1A, 0x2E))

add_text_box(slide, margin_left, margin_top + Inches(0.55), content_width, Inches(0.35),
             "\u57fa\u4e8e\u5b9e\u6d4b\u6570\u636e (2026-03-10, Windows 11, \u4e0a\u6d77\u7f51\u7edc\u73af\u5883)",
             font_size=11, color=RGBColor(0x88, 0x88, 0x88))

# ── Table 1: Our Benchmarks ─────────────────────────────────

add_text_box(slide, margin_left, margin_top + Inches(0.95), Inches(4), Inches(0.3),
             "\u5b9e\u6d4b TTFT \u5bf9\u6bd4 (3\u8f6e\u5e73\u5747)",
             font_size=14, bold=True, color=RGBColor(0x1A, 0x1A, 0x2E))

t1_top = margin_top + Inches(1.3)
t1_shape = slide.shapes.add_table(6, 5, Emu(margin_left), Emu(t1_top), Inches(7.2), Inches(2.4))
t1 = t1_shape.table

col_widths = [Inches(2.4), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2)]
for i, w in enumerate(col_widths):
    t1.columns[i].width = Emu(w)

HEADER_BG = "1A1A2E"
headers = ["\u6a21\u578b / \u65b9\u6848", "TTFT \u5747\u503c", "TTFT \u6700\u5feb", "\u603b\u8017\u65f6", "\u63d0\u5347"]
for i, h in enumerate(headers):
    style_cell(t1.cell(0, i), h, font_size=11, bold=True,
               color=RGBColor(0xFF, 0xFF, 0xFF), bg_color=HEADER_BG)

ROW_LIGHT = "F8F9FA"
ROW_GREEN = "E8F5E9"
ROW_HIGHLIGHT = "C8E6C9"

data = [
    ("Claude Opus 4.6 (relay)", "3.95s", "3.49s", "4.71s", "\u57fa\u7ebf", ROW_LIGHT),
    ("Claude Sonnet 4.5 (relay)", "4.24s", "3.23s", "4.63s", "-", ROW_LIGHT),
    ("Claude Haiku 4.5 (relay)", "4.62s", "3.57s", "5.02s", "-", ROW_LIGHT),
    ("MiniMax M2.5 Lightning", "1.61s", "0.59s", "2.79s", "2.5x", ROW_GREEN),
    ("MiniMax M2.5 + TTS \u5168\u94fe\u8def", "~2.1s", "~1.1s", "~3.3s", "4x", ROW_HIGHLIGHT),
]

for row_idx, (model, ttft_avg, ttft_min, total, improve, bg) in enumerate(data):
    r = row_idx + 1
    is_mm = r >= 4
    style_cell(t1.cell(r, 0), model, font_size=10, bold=is_mm,
               alignment=PP_ALIGN.LEFT, bg_color=bg)
    style_cell(t1.cell(r, 1), ttft_avg, font_size=11, bold=is_mm, bg_color=bg)
    style_cell(t1.cell(r, 2), ttft_min, font_size=11, bold=is_mm, bg_color=bg)
    style_cell(t1.cell(r, 3), total, font_size=11, bold=is_mm, bg_color=bg)
    imp_color = RGBColor(0x2E, 0x7D, 0x32) if is_mm else RGBColor(0x66, 0x66, 0x66)
    style_cell(t1.cell(r, 4), improve, font_size=11, bold=is_mm,
               color=imp_color, bg_color=bg)

# ── Table 2: Industry Benchmarks ─────────────────────────────

t2_label_top = margin_top + Inches(3.85)
add_text_box(slide, margin_left, t2_label_top, Inches(5), Inches(0.3),
             "\u884c\u4e1a\u901f\u5ea6\u6392\u884c (Artificial Analysis, 2026-03)",
             font_size=14, bold=True, color=RGBColor(0x1A, 0x1A, 0x2E))

t2_top = t2_label_top + Inches(0.35)
t2_shape = slide.shapes.add_table(6, 5, Emu(margin_left), Emu(t2_top), Inches(7.2), Inches(2.2))
t2 = t2_shape.table

for i, w in enumerate(col_widths):
    t2.columns[i].width = Emu(w)

headers2 = ["\u6a21\u578b", "TTFT", "\u8f93\u51fa\u901f\u5ea6", "TTS \u5ef6\u8fdf", "\u8bed\u97f3\u603b\u5ef6\u8fdf"]
for i, h in enumerate(headers2):
    style_cell(t2.cell(0, i), h, font_size=11, bold=True,
               color=RGBColor(0xFF, 0xFF, 0xFF), bg_color=HEADER_BG)

industry_data = [
    ("Gemini 2.5 Flash (Google)", "0.4s", "221 TPS", "+0.2s TTS", "~1.0s", ROW_LIGHT),
    ("Claude Haiku 4.5 (\u76f4\u8fde)", "0.6s", "135 TPS", "+0.2s TTS", "~1.2s", ROW_LIGHT),
    ("GPT-4.1 Mini (OpenAI)", "0.5s", "127 TPS", "+0.2s TTS", "~1.1s", ROW_LIGHT),
    ("MiniMax M2.5 Lightning", "1.6s", "100 TPS", "+0.5s TTS", "~2.5s", ROW_GREEN),
    ("Claude via Relay (\u5f53\u524d)", "3.5-5s", "~15 TPS", "+0.5s TTS", "~5-8s", "FFF3E0"),
]

for row_idx, (model, ttft, tps, tts, total, bg) in enumerate(industry_data):
    r = row_idx + 1
    is_mm = r == 4
    is_relay = r == 5
    style_cell(t2.cell(r, 0), model, font_size=10, bold=is_mm,
               alignment=PP_ALIGN.LEFT, bg_color=bg)
    style_cell(t2.cell(r, 1), ttft, font_size=11, bold=is_mm, bg_color=bg)
    style_cell(t2.cell(r, 2), tps, font_size=11, bold=is_mm, bg_color=bg)
    style_cell(t2.cell(r, 3), tts, font_size=11, bg_color=bg)
    if is_mm:
        tc = RGBColor(0x2E, 0x7D, 0x32)
    elif is_relay:
        tc = RGBColor(0xE6, 0x51, 0x00)
    else:
        tc = RGBColor(0x33, 0x33, 0x33)
    style_cell(t2.cell(r, 4), total, font_size=11, bold=True, color=tc, bg_color=bg)

# ── Right panel: Key insights ────────────────────────────────

right_x = margin_left + Inches(7.6)
right_w = Inches(4.8)

add_text_box(slide, right_x, margin_top + Inches(0.95), right_w, Inches(0.3),
             "\u5173\u952e\u7ed3\u8bba",
             font_size=14, bold=True, color=RGBColor(0x1A, 0x1A, 0x2E))

insights = [
    ("\u2705 MiniMax \u5168\u94fe\u8def",
     "LLM 1.6s + TTS 0.5s \u2248 2.1s\n\u6bd4 Relay \u65b9\u6848 (5-8s) \u5feb 3-4 \u500d\nSSE \u6d41\u5f0f TTS \u589e\u91cf\u64ad\u653e\uff0c\u4e0d\u7b49\u5b8c\u6574\u97f3\u9891"),
    ("\u26a0\ufe0f Relay \u74f6\u9888",
     "\u4e2d\u8f6c\u4ee3\u7406\u56fa\u5b9a\u5f00\u9500 ~2-3s\nHaiku \u76f4\u8fde 0.6s \u2192 \u7ecf relay \u53d8 4.6s\n\u975e\u6a21\u578b\u6162\uff0c\u662f relay \u6162"),
    ("\U0001f3c6 \u884c\u4e1a\u6700\u5feb\u65b9\u6848",
     "Gemini Flash / Haiku \u76f4\u8fde + MiniMax TTS\nTTFT 0.4-0.6s + TTS 0.2s \u2248 1s\n\u9700\u8981\u76f4\u8fde API\uff0c\u4e0d\u8d70 relay"),
    ("\U0001f4cb \u5f53\u524d\u63a8\u8350",
     "\u77ed\u671f: MiniMax M2.5 \u5168\u5bb6\u6876 (~2.5s)\n\u4e2d\u671f: Claude Haiku \u76f4\u8fde + MiniMax TTS\n\u957f\u671f: \u7aef\u4fa7 SLM + \u672c\u5730 TTS (\u96f6\u7f51\u7edc)"),
]

y = margin_top + Inches(1.35)
for title, body in insights:
    add_text_box(slide, right_x, y, right_w, Inches(0.25),
                 title, font_size=12, bold=True, color=RGBColor(0x2E, 0x7D, 0x32))
    y += Inches(0.28)
    add_text_box(slide, right_x + Inches(0.15), y, right_w - Inches(0.15), Inches(0.7),
                 body, font_size=10, color=RGBColor(0x55, 0x55, 0x55))
    y += Inches(0.78)

# ── Bottom: TTS bar ──────────────────────────────────────────

add_text_box(slide, margin_left, margin_top + Inches(6.45), content_width, Inches(0.35),
             "TTS \u5b9e\u6d4b: MiniMax Speech 2.8 Turbo TTFT 0.46s \u00b7 7 chunks \u589e\u91cf\u64ad\u653e \u00b7 9.5s \u97f3\u9891\u4ec5 0.79s \u751f\u6210  |  Edge TTS ~3s (\u5168\u91cf\u7f13\u51b2)  |  \u672c\u5730 sherpa-onnx ~0.5s (\u79bb\u7ebf)",
             font_size=10, color=RGBColor(0x66, 0x66, 0x66))

# ── Save ─────────────────────────────────────────────────────

output = "docs/NOVA_Dog_Askme_v4_Report_v6.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")

# Verify slide order
for i, sl in enumerate(prs.slides):
    texts = []
    for shape in sl.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
                    break
        if texts:
            break
    print(f"  Slide {i+1}: {texts[0][:50] if texts else '(empty)'}")
