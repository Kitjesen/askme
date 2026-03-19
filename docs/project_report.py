"""NOVA Dog Askme — Project Report PPT (v3 rewrite: innovation-focused, light theme)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (light, clean) ────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG         = RGBColor(0xFA, 0xFA, 0xFC)
CARD       = RGBColor(0xF0, 0xF2, 0xF5)
TITLE      = RGBColor(0x11, 0x18, 0x27)
BODY       = RGBColor(0x33, 0x40, 0x55)
SUB        = RGBColor(0x64, 0x74, 0x8B)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)

BLUE       = RGBColor(0x25, 0x63, 0xEB)
BLUE_BG    = RGBColor(0xEF, 0xF6, 0xFF)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_BG  = RGBColor(0xF5, 0xF3, 0xFF)
GREEN      = RGBColor(0x05, 0x96, 0x69)
GREEN_BG   = RGBColor(0xEC, 0xFD, 0xF5)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_BG  = RGBColor(0xFF, 0xF7, 0xED)
RED        = RGBColor(0xDC, 0x26, 0x26)
RED_BG     = RGBColor(0xFE, 0xF2, 0xF2)
CYAN       = RGBColor(0x06, 0x84, 0x8D)
CYAN_BG    = RGBColor(0xF0, 0xFD, 0xFA)

FONT = "Microsoft YaHei"
W = 13.333  # slide width
H = 7.5     # slide height


# ── Helpers ────────────────────────────────────────────────────
def _bg(slide, color=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _txt(slide, l, t, w, h, text, sz=16, c=BODY, bold=False, al=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = FONT; p.alignment = al
    return tf

def _p(tf, text, sz=14, c=BODY, bold=False, sp=Pt(4)):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = FONT; p.space_before = sp
    return p

def _card(slide, l, t, w, h, fill=CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False

def _bar(slide, l, t, w, c=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def _header(slide, text, sub=None):
    """Standard slide header with accent bar + title."""
    _bar(slide, 0.7, 0.4, 0.5, BLUE)
    _txt(slide, 1.35, 0.28, 10, 0.6, text, sz=28, c=TITLE, bold=True)
    if sub:
        _txt(slide, 1.35, 0.78, 10, 0.4, sub, sz=15, c=SUB)


def build(out):
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    B = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════
    # SLIDE 1 — 封面: 一句话讲清楚
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _bar(s, 0, 0, W, BLUE)

    _txt(s, 1.2, 1.0, 10, 1.0, "Askme", sz=52, c=TITLE, bold=True)
    _txt(s, 1.2, 2.0, 10, 0.6,
         "让工业机器人听懂人话、记住经验、安全执行",
         sz=26, c=BLUE, bold=True)
    _bar(s, 1.2, 2.8, 4, BLUE)
    _txt(s, 1.2, 3.2, 9, 0.8,
         "一套面向工业机器人的语音 AI 运行时\n"
         "不是聊天机器人包装 — 是让 LLM 变成可信赖的现场 Agent",
         sz=17, c=BODY)

    _txt(s, 1.2, 4.5, 6, 0.4, "上海穹沛科技  ·  2026年3月", sz=14, c=MUTED)

    # Right side: 3 core numbers that matter
    highlights = [
        ("5 大创新", "区别于所有现有方案", BLUE, BLUE_BG),
        ("4 层记忆", "机器人会学习、会反思", PURPLE, PURPLE_BG),
        ("3 级安全", "危险操作必须人确认", RED, RED_BG),
    ]
    for i, (num, desc, clr, bg) in enumerate(highlights):
        y = 1.2 + i * 1.2
        _card(s, 9.0, y, 3.8, 0.95, fill=bg)
        _txt(s, 9.3, y + 0.1, 2.0, 0.4, num, sz=22, c=clr, bold=True)
        _txt(s, 9.3, y + 0.5, 3.2, 0.35, desc, sz=13, c=SUB)

    # ══════════════════════════════════════════════════════════
    # SLIDE 2 — 解决什么问题
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "工业机器人语音交互的 6 个痛点")

    problems = [
        ("延迟高", "商业助手 2-5s 往返\n机器人像在发呆", "流式 TTS\n边生成边播放", ORANGE, ORANGE_BG),
        ("不安全", "误听一句就执行危险动作\n没有确认机制", "三级审批\n危险操作必须人确认", RED, RED_BG),
        ("会忘事", "对话结束就忘了\n每次从零开始", "四层记忆\n经验 → 反思 → 知识", PURPLE, PURPLE_BG),
        ("不可控", "云端黑盒\n网络断了就哑了", "本地 ASR/VAD/TTS\n离线也能听说", BLUE, BLUE_BG),
        ("不可审", "谁下的令？何时执行？\n无法追溯", "全链路事件溯源\n操作员身份追踪", CYAN, CYAN_BG),
        ("不可复用", "每个机器人重写一遍\n技能无法共享", "技能包插拔\n新场景零代码扩展", GREEN, GREEN_BG),
    ]
    for i, (pain, detail, solution, clr, bg) in enumerate(problems):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 1.4 + row * 2.9

        _card(s, x, y, 3.9, 2.6, fill=bg)
        _txt(s, x + 0.2, y + 0.15, 3.5, 0.35, f"痛点: {pain}", sz=15, c=clr, bold=True)
        lines = detail.split("\n")
        _txt(s, x + 0.2, y + 0.55, 3.5, 0.3, lines[0], sz=12, c=SUB)
        if len(lines) > 1:
            _txt(s, x + 0.2, y + 0.8, 3.5, 0.3, lines[1], sz=12, c=SUB)

        _bar(s, x + 0.2, y + 1.15, 1.5, clr)
        _txt(s, x + 0.2, y + 1.3, 3.5, 0.35, "Askme 方案:", sz=13, c=TITLE, bold=True)
        sol_lines = solution.split("\n")
        _txt(s, x + 0.2, y + 1.65, 3.5, 0.3, sol_lines[0], sz=13, c=clr, bold=True)
        if len(sol_lines) > 1:
            _txt(s, x + 0.2, y + 1.95, 3.5, 0.3, sol_lines[1], sz=12, c=BODY)

    # ══════════════════════════════════════════════════════════
    # SLIDE 3 — 创新 1: 四层记忆
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "创新 ①  四层记忆 — 机器人会学习、会反思",
            "不是 RAG，不是微调，是认知架构")

    layers = [
        ("L1 对话", "滑动窗口", BLUE, BLUE_BG,
         ["最近 40 条消息", "超量自动 LLM 压缩为摘要", "实时持久化 JSON"]),
        ("L2 会话", "中期摘要", PURPLE, PURPLE_BG,
         ["对话结束 → 2-4句摘要", "按日期存储 .md 文件", "下次启动注入上下文"]),
        ("L3 情景", "反思+知识", GREEN, GREEN_BG,
         ["事件日志 + 重要性评分", "Ebbinghaus 遗忘曲线", "自动反思 → 5 类世界知识"]),
        ("L4 向量", "长期语义", CYAN, CYAN_BG,
         ["Qwen3 Embedding 0.6B", "语义检索 (RAG)", "懒加载，服务断不崩溃"]),
    ]
    x = 0.3
    for title, subtitle, clr, bg, items in layers:
        _card(s, x, 1.3, 3.1, 3.2, fill=bg)
        _txt(s, x + 0.2, 1.4, 2.7, 0.3, title, sz=18, c=clr, bold=True)
        _txt(s, x + 0.2, 1.75, 2.7, 0.25, subtitle, sz=12, c=SUB)
        y = 2.1
        for item in items:
            _txt(s, x + 0.2, y, 2.7, 0.25, f"• {item}", sz=12, c=BODY)
            y += 0.35
        x += 3.25

    # Aha moment
    _card(s, 0.5, 4.8, 12.3, 2.0, fill=PURPLE_BG)
    _txt(s, 0.8, 4.95, 11.5, 0.4,
         "\"Aha\" 时刻: 机器人每天巡检，第 7 天开始变聪明", sz=18, c=PURPLE, bold=True)
    tf = _txt(s, 0.8, 5.45, 5.5, 0.3, "", sz=13, c=BODY)
    _p(tf, "Day 1: 通用问答，不认识任何东西", sz=13, c=SUB)
    _p(tf, "Day 3: 记住了哪些门容易卡（情景记忆）", sz=13, c=BODY)
    _p(tf, "Day 5: 识别出固定时间出现的人员（模式反思）", sz=13, c=BODY)
    _p(tf, "Day 7: 发现设备异常规律，主动预警（世界知识）", sz=13, c=GREEN, bold=True)

    _txt(s, 7.0, 5.45, 5.5, 0.35,
         "L3 情景记忆核心机制:", sz=14, c=TITLE, bold=True)
    tf = _txt(s, 7.0, 5.85, 5.5, 0.3, "", sz=12, c=BODY)
    _p(tf, "重要性评分: command(0.7) + 检测人(+0.5) + 危险(+0.6)", sz=12, c=BODY)
    _p(tf, "遗忘曲线: R = e^(-t/S), 每次回忆 S×2 (间隔效应)", sz=12, c=BODY)
    _p(tf, "反思触发: 累积重要性 ≥ 15 → LLM 提炼知识", sz=12, c=BODY)
    _p(tf, "5 类知识: 环境 · 实体 · 规律 · 交互 · 自省", sz=12, c=PURPLE)

    # ══════════════════════════════════════════════════════════
    # SLIDE 4 — 创新 2: 安全第一
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "创新 ②  安全第一 — ESTOP 零延迟 + 三级审批",
            "工业机器人的安全不是 feature，是底线")

    # Intent routing
    _card(s, 0.5, 1.3, 6.0, 2.5, fill=RED_BG)
    _txt(s, 0.8, 1.4, 5.4, 0.35,
         "意图路由 — 安全优先级", sz=16, c=RED, bold=True)
    intents = [
        ("① ESTOP", "\"停！\" \"紧急停止\"", "直达 safety 服务，不经 LLM", RED),
        ("② COMMAND", "\"坐下\" \"站起来\"", "确定性指令，不需要推理", ORANGE),
        ("③ TRIGGER", "\"去厨房\" \"开始巡检\"", "匹配 Skill，走任务编排", BLUE),
        ("④ GENERAL", "\"你好\" \"今天天气\"", "才送到 LLM 自由对话", SUB),
    ]
    y = 1.85
    for name, example, desc, clr in intents:
        _txt(s, 0.8, y, 1.5, 0.25, name, sz=13, c=clr, bold=True)
        _txt(s, 2.3, y, 2.0, 0.25, example, sz=12, c=BODY)
        _txt(s, 4.2, y, 2.0, 0.25, desc, sz=11, c=SUB)
        y += 0.4

    # Tool safety levels
    _card(s, 6.8, 1.3, 6.0, 2.5, fill=ORANGE_BG)
    _txt(s, 7.1, 1.4, 5.4, 0.35,
         "工具三级安全", sz=16, c=ORANGE, bold=True)
    levels = [
        ("normal", "直接执行", "查时间、读文件、看状态", GREEN),
        ("dangerous", "需确认 \"确认执行\"", "移动机械臂、网络搜索、执行命令", ORANGE),
        ("critical", "最高警戒", "紧急停止 (bypass: 免审批直达)", RED),
    ]
    y = 1.85
    for level, action, examples, clr in levels:
        _txt(s, 7.1, y, 1.3, 0.25, level, sz=13, c=clr, bold=True)
        _txt(s, 8.4, y, 1.8, 0.25, action, sz=12, c=BODY)
        _txt(s, 10.2, y, 2.4, 0.25, examples, sz=11, c=SUB)
        y += 0.4

    # Dialogue example
    _card(s, 0.5, 4.1, 12.3, 3.0, fill=CARD)
    _txt(s, 0.8, 4.2, 5, 0.35,
         "真实对话场景 — 安全审批", sz=16, c=TITLE, bold=True)

    dialogue = [
        ("操作员:", "\"把机械臂移到前面\"", BLUE),
        ("Askme:", "识别为 robot_move (dangerous) → 需要确认", SUB),
        ("Thunder:", "\"这是高风险操作: 移动机械臂到 (0, 100, 0)。请说'确认执行'继续，或'取消'放弃。\"", ORANGE),
        ("操作员:", "\"确认执行\"", BLUE),
        ("Thunder:", "\"收到，正在移动机械臂...已到达目标位置。\"", GREEN),
    ]
    y = 4.65
    for speaker, text, clr in dialogue:
        _txt(s, 0.8, y, 1.2, 0.3, speaker, sz=13, c=clr, bold=True)
        _txt(s, 2.0, y, 10.5, 0.3, text, sz=13, c=BODY)
        y += 0.42

    _txt(s, 0.8, 6.5, 12, 0.3,
         "审批超时 30s 自动拒绝 · 确认词: 确认执行/继续执行/approve · 拒绝词: 取消/cancel",
         sz=12, c=MUTED)

    # ══════════════════════════════════════════════════════════
    # SLIDE 5 — 创新 3: 流式语音 + 创新 4: 插拔技能
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "创新 ③ 流式语音  ·  创新 ④ 插拔技能")

    # Streaming TTS
    _card(s, 0.5, 1.2, 6.0, 3.0, fill=BLUE_BG)
    _txt(s, 0.8, 1.3, 5.4, 0.35, "流式 TTS — 不等全部生成完就开口",
         sz=16, c=BLUE, bold=True)
    tf = _txt(s, 0.8, 1.75, 5.4, 0.3, "", sz=13, c=BODY)
    _p(tf, "传统方案: LLM 全部输出 → TTS 合成 → 播放 (等 8s+)", sz=13, c=SUB)
    _p(tf, "Askme: LLM 每出一句 → 立即 TTS → 立即播放", sz=13, c=BLUE, bold=True)
    _p(tf, "", sz=6)
    _p(tf, "StreamSplitter 5 规则智能分句:", sz=13, c=TITLE, bold=True)
    _p(tf, "  。？！→ 立即分割 (强标点)", sz=12, c=BODY)
    _p(tf, "  ；：  → 8字后分割 (中等标点)", sz=12, c=BODY)
    _p(tf, "  ，    → 15字后分割 (逗号)", sz=12, c=BODY)
    _p(tf, "  60字  → 强制输出 (防卡顿)", sz=12, c=BODY)
    _p(tf, "", sz=4)
    _p(tf, "效果: 首句音频在 LLM 第一个句号后 ~500ms 开始播放", sz=13, c=GREEN, bold=True)

    # Skills system
    _card(s, 6.8, 1.2, 6.0, 3.0, fill=GREEN_BG)
    _txt(s, 7.1, 1.3, 5.4, 0.35, "Skill = 一个 Markdown 文件",
         sz=16, c=GREEN, bold=True)
    tf = _txt(s, 7.1, 1.75, 5.4, 0.3, "", sz=13, c=BODY)
    _p(tf, "添加新技能 = 定义一个技能描述文件", sz=13, c=GREEN, bold=True)
    _p(tf, "无需改代码 · 无需重新部署 · 即放即用", sz=12, c=SUB)
    _p(tf, "", sz=6)
    _p(tf, "每个技能包含:", sz=13, c=TITLE, bold=True)
    _p(tf, "  触发条件 — 语音关键词 / 自动 / 定时", sz=12, c=BODY)
    _p(tf, "  安全级别 — 普通 / 危险 / 关键", sz=12, c=BODY)
    _p(tf, "  执行逻辑 — AI 提示词 + 可用工具范围", sz=12, c=BODY)
    _p(tf, "", sz=4)
    _p(tf, "三级覆盖: 内置 → 用户自定义 → 项目专用", sz=13, c=BODY)
    _p(tf, "后加载覆盖前 — 客户可定制内置技能", sz=12, c=SUB)

    # Skill examples
    _txt(s, 0.5, 4.5, 8, 0.4, "12 个内置技能 — 开箱即用", sz=16, c=TITLE, bold=True)

    skill_list = [
        ("紧急停止", "robot_estop", "\"紧急停止\"", "critical", RED),
        ("移动控制", "robot_move", "\"移动到\"", "dangerous", ORANGE),
        ("抓取操作", "robot_grab", "\"抓住\"", "dangerous", ORANGE),
        ("环境感知", "environment_report", "\"看看周围\"", "normal", GREEN),
        ("巡逻报告", "patrol_report", "\"巡逻报告\"", "normal", GREEN),
        ("时间查询", "get_time", "\"现在几点\"", "normal", GREEN),
        ("自动问候", "greet_person", "\"打个招呼\"", "normal", BLUE),
        ("网络搜索", "web_search", "\"帮我搜索\"", "dangerous", ORANGE),
    ]
    for i, (cn, en, trigger, level, clr) in enumerate(skill_list):
        col = i % 4
        row = i // 4
        x = 0.5 + col * 3.15
        y = 5.0 + row * 0.85
        _card(s, x, y, 2.95, 0.7, fill=CARD)
        _txt(s, x + 0.15, y + 0.05, 1.0, 0.25, cn, sz=13, c=clr, bold=True)
        _txt(s, x + 1.15, y + 0.05, 1.7, 0.25, trigger, sz=11, c=BODY)
        _txt(s, x + 0.15, y + 0.35, 1.5, 0.25, en, sz=10, c=SUB)
        _txt(s, x + 1.8, y + 0.35, 1.0, 0.25, level, sz=10, c=clr)

    # ══════════════════════════════════════════════════════════
    # SLIDE 6 — 创新 5: 产品赋能
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "创新 ⑤  语音 AI 基础设施 — 一次构建，赋能所有机器人",
            "Askme 不是一个产品，是一个平台")

    # Core proposition: 3 cards
    props = [
        ("硬件团队零 AI 门槛", BLUE, BLUE_BG,
         "机器人公司的核心能力是硬件和运动控制\n"
         "不应该再养一个语音 AI 团队\n\n"
         "接入 Askme = 立即获得:\n"
         "语音交互 + 安全审批 + 记忆学习 + 任务编排"),
        ("场景快速复制", GREEN, GREEN_BG,
         "新客户 ≠ 重新开发\n"
         "新场景 ≠ 写新代码\n\n"
         "通过配置面板调整:\n"
         "机器人人格 · 语言风格 · 行为边界\n"
         "技能包 · 安全策略 · 硬件适配\n\n"
         "从签单到交付: 小时级，不是月级"),
        ("越用越值钱", PURPLE, PURPLE_BG,
         "传统方案: 部署后能力固定\n"
         "Askme: 部署后持续进化\n\n"
         "四层记忆让每台机器人积累经验\n"
         "同一条走廊走 100 遍后的机器人\n"
         "比新出厂的机器人更懂这个环境"),
    ]
    x = 0.5
    for title, clr, bg, desc in props:
        _card(s, x, 1.25, 4.0, 3.5, fill=bg)
        _txt(s, x + 0.2, 1.35, 3.6, 0.35, title, sz=16, c=clr, bold=True)
        lines = desc.split("\n")
        ty = 1.8
        for line in lines:
            if line == "":
                ty += 0.1
                continue
            _txt(s, x + 0.2, ty, 3.6, 0.25, line, sz=12, c=BODY if line[0] != "从" else clr,
                 bold=(line.startswith("从") or line.startswith("Askme") or line.startswith("四层")))
            ty += 0.28
        x += 4.15

    # Bottom: what this means for穹沛
    _card(s, 0.5, 5.05, 12.3, 2.1, fill=CARD)
    _txt(s, 0.8, 5.15, 11.5, 0.35,
         "对穹沛意味着什么", sz=18, c=TITLE, bold=True)

    biz = [
        ("巡检整机方案", "Thunder + Askme + LingTu 打包\n语音控制导航，开箱即用",
         "定制人格 + 巡检技能包", BLUE),
        ("跨场景复制", "同一平台，换技能包即可\n仓储/安防/接待快速复制",
         "新场景交付周期: 天级", GREEN),
        ("赋能集成商", "Askme 作为语音 AI 中间件\n集成商用它赋能自己的机器人",
         "按设备授权，持续收入", PURPLE),
        ("远程运维", "OTA 推送新技能/新策略\n远程升级所有设备的能力",
         "不用派人到现场", ORANGE),
    ]
    x = 0.5
    for title, desc, note, clr in biz:
        _card(s, x + 0.1, 5.6, 2.9, 1.35, fill=WHITE)
        _txt(s, x + 0.25, 5.65, 2.6, 0.3, title, sz=13, c=clr, bold=True)
        lines = desc.split("\n")
        _txt(s, x + 0.25, 5.95, 2.6, 0.25, lines[0], sz=11, c=BODY)
        if len(lines) > 1:
            _txt(s, x + 0.25, 6.2, 2.6, 0.25, lines[1], sz=11, c=BODY)
        _txt(s, x + 0.25, 6.5, 2.6, 0.25, note, sz=10, c=SUB)
        x += 3.1

    # ══════════════════════════════════════════════════════════
    # SLIDE 7 — 全链路: 从说话到执行
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "全链路 — \"小穹，去厨房巡检\"")

    # Simplified 3-row flow
    row1 = [
        ("用户说话", "\"小穹，去厨房巡检\"", BLUE),
        ("唤醒 + 识别", "KWS → VAD → ASR\n\"去厨房巡检\"", PURPLE),
        ("意图路由", "VOICE_TRIGGER\nskill: patrol", ORANGE),
        ("边缘服务", "askme-edge\n创建任务", BLUE),
    ]
    row2 = [
        ("安全预检", "dog-safety\nestop? 电量? 黑名单?", RED),
        ("任务编排", "arbiter\nrequested → approved", PURPLE),
        ("导航执行", "nav-gateway\n→ LingTu gRPC", GREEN),
        ("遥测记录", "telemetry-hub\n事件 + 审计日志", CYAN),
    ]
    row3 = [
        ("语音回复", "TTS: \"好的，正在去\n厨房巡检\"", BLUE),
        ("记忆更新", "L3: log(command)\nL1: 加入对话历史", PURPLE),
        ("后台反思", "情景缓冲 → 反思\n→ 世界知识更新", GREEN),
        ("持续学习", "下次再去厨房时\n已经记住路线", CYAN),
    ]

    for row_idx, row in enumerate([row1, row2, row3]):
        y = 1.2 + row_idx * 2.0
        for i, (title, desc, clr) in enumerate(row):
            x = 0.4 + i * 3.2
            bg = {BLUE: BLUE_BG, PURPLE: PURPLE_BG, ORANGE: ORANGE_BG,
                  RED: RED_BG, GREEN: GREEN_BG, CYAN: CYAN_BG}[clr]
            _card(s, x, y, 2.8, 1.6, fill=bg)
            _txt(s, x + 0.15, y + 0.1, 2.5, 0.3, title, sz=14, c=clr, bold=True)
            lines = desc.split("\n")
            _txt(s, x + 0.15, y + 0.5, 2.5, 0.3, lines[0], sz=12, c=BODY)
            if len(lines) > 1:
                _txt(s, x + 0.15, y + 0.8, 2.5, 0.3, lines[1], sz=12, c=SUB)
            # Arrow
            if i < 3:
                _txt(s, x + 2.7, y + 0.5, 0.5, 0.4, "→", sz=18, c=clr, bold=True)

        # Down arrow between rows
        if row_idx < 2:
            _txt(s, 12.2, y + 1.3, 0.5, 0.4, "↓", sz=18, c=MUTED, bold=True)

    # ══════════════════════════════════════════════════════════
    # SLIDE 8 — 产品方案总览
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "巡检导航完整方案 — Askme + LingTu + NOVA Runtime",
            "语音操控 + 自主导航 + 安全审批 + 持续学习 = 开箱即用")

    prod_mods = [
        ("Askme 语音 AI", "语音交互运行时",
         ["唤醒 → 识别 → 意图 → LLM 对话",
          "流式 TTS 即时播报",
          "三级安全审批",
          "四层记忆持续学习",
          "12+ 内置技能"],
         BLUE, BLUE_BG),
        ("LingTu 导航", "v1.7.5 稳定运行",
         ["SLAM 建图",
          "路径规划 + 避障",
          "6 种导航模式",
          "多楼层支持",
          "gRPC 标准接口"],
         GREEN, GREEN_BG),
        ("NOVA Runtime", "6 微服务编排",
         ["任务状态机",
          "安全守卫 ESTOP",
          "遥测 + 审计",
          "Docker 一键部署",
          "共享认证"],
         PURPLE, PURPLE_BG),
        ("OTA 运维", "空中升级",
         ["远程推送新技能",
          "安全策略更新",
          "设备健康监控",
          "固件升级",
          "不用派人到现场"],
         ORANGE, ORANGE_BG),
    ]
    x = 0.3
    for title, sub, items, clr, bg in prod_mods:
        _card(s, x, 1.2, 3.1, 3.5, fill=bg)
        _txt(s, x + 0.15, 1.3, 2.8, 0.3, title, sz=15, c=clr, bold=True)
        _txt(s, x + 0.15, 1.6, 2.8, 0.25, sub, sz=11, c=SUB)
        _bar(s, x + 0.15, 1.9, 1.0, clr)
        ty = 2.05
        for item in items:
            _txt(s, x + 0.15, ty, 2.8, 0.25, f"• {item}", sz=11, c=BODY)
            ty += 0.28
        x += 3.2

    # Use cases
    scenes = [
        ("电力巡检", "变电站·配电房·线路", BLUE),
        ("化工安防", "管廊·罐区·危险区域", RED),
        ("仓储物流", "货架·通道·月台", GREEN),
        ("园区安保", "楼宇·停车场·周界", PURPLE),
    ]
    x = 0.3
    for title, area, clr in scenes:
        bg = {BLUE: BLUE_BG, RED: RED_BG, GREEN: GREEN_BG, PURPLE: PURPLE_BG}[clr]
        _card(s, x, 5.0, 3.1, 0.9, fill=bg)
        _txt(s, x + 0.15, 5.05, 1.2, 0.25, title, sz=13, c=clr, bold=True)
        _txt(s, x + 1.3, 5.05, 1.7, 0.25, area, sz=10, c=SUB)
        x += 3.2

    # Operator capabilities
    _bar(s, 0.7, 6.1, 0.4, BLUE)
    _txt(s, 1.2, 6.0, 4, 0.4, "操作员语音指令", sz=16, c=TITLE, bold=True)
    caps = [
        ("\"去三号柜\"", "巡检", GREEN),
        ("\"电池多少\"", "查询", BLUE),
        ("\"拍照记录\"", "感知", PURPLE),
        ("\"停！\"", "急停", RED),
        ("\"取消任务\"", "管理", ORANGE),
        ("\"上次这里?\"", "记忆", CYAN),
    ]
    x = 0.3
    for cmd, label, clr in caps:
        _card(s, x, 6.4, 2.05, 0.7, fill={GREEN:GREEN_BG,BLUE:BLUE_BG,PURPLE:PURPLE_BG,RED:RED_BG,ORANGE:ORANGE_BG,CYAN:CYAN_BG}[clr])
        _txt(s, x + 0.1, 6.42, 1.0, 0.25, cmd, sz=11, c=BODY)
        _txt(s, x + 1.1, 6.42, 0.8, 0.25, label, sz=10, c=clr, bold=True)
        x += 2.15

    # ══════════════════════════════════════════════════════════
    # SLIDE 9 — 语音效果演示 (上)
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "实际语音效果 — 本地离线 TTS 真实音频",
            "VITS-Aishell3 · RTF 0.2 (5x 实时) · 纯 CPU · 329MB · 零网络依赖")

    demo_up = [
        ("① 唤醒+问候", "操作员:", "\"小穹\"  (唤醒词触发)",
         "Thunder:", "你好，我是Thunder，穹沛的巡检机器人。等待指令。",
         "01_greeting.wav  3.7s", "KWS → ASR", BLUE, BLUE_BG),

        ("② 下达巡检", "操作员:", "\"去三号配电柜巡检\"",
         "Thunder:", "收到，正在前往三号配电柜巡检。预计到达时间两分钟。",
         "02_confirm_task.wav  5.4s", "Intent → Skill → nav-gateway", GREEN, GREEN_BG),

        ("③ 安全审批", "操作员:", "\"把机械臂移到前面\"  (dangerous 工具)",
         "Thunder:", "这是高风险操作。请说'确认执行'继续，或'取消'放弃。",
         "04_safety_confirm.wav  6.9s", "Safety 三级审批", ORANGE, ORANGE_BG),
    ]

    y = 1.2
    for step, sp1, t1, sp2, t2, wav_info, sys_tag, clr, bg in demo_up:
        _card(s, 0.3, y, 12.7, 1.65, fill=bg)
        _txt(s, 0.5, y + 0.05, 2.5, 0.3, step, sz=15, c=clr, bold=True)
        _txt(s, 7.5, y + 0.05, 5.0, 0.25, sys_tag, sz=11, c=clr, al=PP_ALIGN.RIGHT)
        _txt(s, 0.5, y + 0.4, 0.9, 0.25, sp1, sz=12, c=BLUE, bold=True)
        _txt(s, 1.4, y + 0.4, 6.0, 0.25, t1, sz=12, c=BODY)
        _txt(s, 0.5, y + 0.7, 0.9, 0.25, sp2, sz=12, c=clr, bold=True)
        _txt(s, 1.4, y + 0.7, 8.0, 0.45, t2, sz=12, c=TITLE)

        # Audio indicator
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(10.8), Inches(y + 0.3), Inches(0.7), Inches(0.7))
        circ.fill.solid(); circ.fill.fore_color.rgb = clr
        circ.line.fill.background()
        tf = circ.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.text = "▶"; p.font.size = Pt(16)
        p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

        _txt(s, 10.2, y + 1.1, 2.8, 0.25, wav_info, sz=9, c=SUB, al=PP_ALIGN.CENTER)

        y += 1.8

    _card(s, 0.3, 6.65, 12.7, 0.5, fill=CARD)
    _txt(s, 0.5, 6.7, 12.3, 0.35,
         "音频文件位于 docs/voice_samples/ — 演示时直接播放 WAV 或在 PowerPoint 中插入音频",
         sz=12, c=SUB, al=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 10 — 语音效果演示 (下)
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "语音效果 (续) — 异常报告 · 急停 · 记忆调用")

    demo_dn = [
        ("④ 巡检报告", "Thunder:", "(到达目标，自动检测异常)",
         "Thunder:", "三号区域检测到温度异常，柜体表面六十八度，超过阈值。已记录并上报。",
         "05_patrol_report.wav  8.2s", "Skill → Telemetry", PURPLE, PURPLE_BG),

        ("⑤ 紧急停止", "操作员:", "\"停！\"  (ESTOP Tier-1, 不经 LLM)",
         "Thunder:", "已紧急停机。所有运动已暂停，等待下一步指令。",
         "03_estop.wav  4.3s", "ESTOP <100ms", RED, RED_BG),

        ("⑥ 记忆调用", "操作员:", "\"之前这里有什么异常吗？\"",
         "Thunder:", "根据之前巡检记录，这个区域上周也出现过类似温度异常。建议排查散热系统。",
         "06_memory_demo.wav  6.9s", "L3 情景记忆检索", CYAN, CYAN_BG),
    ]

    y = 1.2
    for step, sp1, t1, sp2, t2, wav_info, sys_tag, clr, bg in demo_dn:
        _card(s, 0.3, y, 12.7, 1.65, fill=bg)
        _txt(s, 0.5, y + 0.05, 2.5, 0.3, step, sz=15, c=clr, bold=True)
        _txt(s, 7.5, y + 0.05, 5.0, 0.25, sys_tag, sz=11, c=clr, al=PP_ALIGN.RIGHT)
        _txt(s, 0.5, y + 0.4, 0.9, 0.25, sp1, sz=12,
             c=BLUE if "操作" in sp1 else clr, bold=True)
        _txt(s, 1.4, y + 0.4, 6.0, 0.25, t1, sz=12, c=BODY)
        _txt(s, 0.5, y + 0.7, 0.9, 0.25, sp2, sz=12, c=clr, bold=True)
        _txt(s, 1.4, y + 0.7, 8.0, 0.45, t2, sz=12, c=TITLE)

        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(10.8), Inches(y + 0.3), Inches(0.7), Inches(0.7))
        circ.fill.solid(); circ.fill.fore_color.rgb = clr
        circ.line.fill.background()
        tf = circ.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.text = "▶"; p.font.size = Pt(16)
        p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

        _txt(s, 10.2, y + 1.1, 2.8, 0.25, wav_info, sz=9, c=SUB, al=PP_ALIGN.CENTER)

        y += 1.8

    # TTS performance summary
    _card(s, 0.3, 6.4, 6.0, 0.8, fill=GREEN_BG)
    _txt(s, 0.5, 6.45, 5.6, 0.3, "TTS 合成性能 (实测)", sz=14, c=GREEN, bold=True)
    _txt(s, 0.5, 6.75, 5.6, 0.25,
         "短句 0.9s · 长句 1.8s · RTF=0.2 · CPU · 329MB 模型 · 4 种音色可选",
         sz=11, c=BODY)

    _card(s, 6.6, 6.4, 6.4, 0.8, fill=BLUE_BG)
    _txt(s, 6.8, 6.45, 6.0, 0.3, "ASR 识别性能", sz=14, c=BLUE, bold=True)
    _txt(s, 6.8, 6.75, 6.0, 0.25,
         "中英双语 · 流式识别 · 530MB 模型 · 离线可用 · VAD 过滤噪声",
         sz=11, c=BODY)

    # ══════════════════════════════════════════════════════════
    # SLIDE 11 — Runtime 服务 + 模块化 (原 SLIDE 8)
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "6 个微服务 — 每个都可独立部署",
            "Docker Compose 一键启动 · 共享认证 · 事件溯源")

    svcs = [
        ("askme-edge", ":5100", "语音入口",
         "接收语音 → 路由 → 聚合状态 → 返回回复", BLUE, BLUE_BG),
        ("mission", ":5050", "任务编排",
         "全生命周期状态机 · 安全预检 · 审批流 · 审计", PURPLE, PURPLE_BG),
        ("safety", ":5070", "安全守卫",
         "ESTOP · 维护模式 · 电量阈值 · 能力黑名单", RED, RED_BG),
        ("control", ":5080", "动作执行",
         "8 种 capability → Brainstem gRPC (等硬件)", ORANGE, ORANGE_BG),
        ("nav", ":5090", "导航网关",
         "6 种模式 → LingTu gRPC (可立即接入)", GREEN, GREEN_BG),
        ("telemetry", ":5060", "遥测中心",
         "事件流 · 审计日志 · 跨服务可观测", CYAN, CYAN_BG),
    ]
    for i, (name, port, role, desc, clr, bg) in enumerate(svcs):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = 1.3 + row * 1.45
        _card(s, x, y, 6.1, 1.25, fill=bg)
        _txt(s, x + 0.2, y + 0.1, 2.5, 0.3, f"{name} {port}", sz=15, c=clr, bold=True)
        _txt(s, x + 2.7, y + 0.1, 1.5, 0.3, role, sz=13, c=TITLE, bold=True)
        _txt(s, x + 0.2, y + 0.5, 5.7, 0.6, desc, sz=12, c=BODY)

    # Modular pattern
    _bar(s, 0.7, 5.8, 0.4, GREEN)
    _txt(s, 1.2, 5.7, 8, 0.4, "添加新服务只需 3 步", sz=18, c=TITLE, bold=True)

    steps = [
        ("1. 写服务", "标准模板:\napp.py + api.py + service.py", GREEN, GREEN_BG),
        ("2. 加认证", "NOVA_DOG_RUNTIME_API_KEY\n共享一把钥匙", BLUE, BLUE_BG),
        ("3. 部署", "docker-compose.yml\n加一段就上线", PURPLE, PURPLE_BG),
    ]
    x = 0.5
    for title, desc, clr, bg in steps:
        _card(s, x, 6.15, 4.0, 1.0, fill=bg)
        _txt(s, x + 0.2, 6.2, 1.2, 0.3, title, sz=14, c=clr, bold=True)
        lines = desc.split("\n")
        _txt(s, x + 1.4, 6.2, 2.5, 0.3, lines[0], sz=12, c=BODY)
        if len(lines) > 1:
            _txt(s, x + 1.4, 6.5, 2.5, 0.3, lines[1], sz=11, c=SUB)
        if x < 8:
            _txt(s, x + 3.9, 6.4, 0.3, 0.3, "→", sz=16, c=clr, bold=True)
        x += 4.2

    # ══════════════════════════════════════════════════════════
    # SLIDE 9 — 延迟优化
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "延迟 — 怎么让机器人不像在发呆")

    # Before vs after
    _card(s, 0.5, 1.2, 6.0, 2.3, fill=RED_BG)
    _txt(s, 0.8, 1.3, 5.4, 0.35, "传统方案 (串行)", sz=16, c=RED, bold=True)
    tf = _txt(s, 0.8, 1.75, 5.4, 0.3, "", sz=13, c=BODY)
    _p(tf, "ASR (3s) → 记忆检索 (3s) → LLM (6s) → TTS (2s)", sz=13, c=BODY)
    _p(tf, "总计: ~14 秒才开口说话", sz=15, c=RED, bold=True)
    _p(tf, "用户体验: 机器人像死机了一样", sz=13, c=SUB)

    _card(s, 6.8, 1.2, 6.0, 2.3, fill=GREEN_BG)
    _txt(s, 7.1, 1.3, 5.4, 0.35, "Askme 方案 (并行 + 流式)", sz=16, c=GREEN, bold=True)
    tf = _txt(s, 7.1, 1.75, 5.4, 0.3, "", sz=13, c=BODY)
    _p(tf, "ASR + 记忆预取 (并行 3s) → LLM 流式 → 边说边播", sz=13, c=BODY)
    _p(tf, "首句音频: ~5-6 秒 (Opus)  /  ~3 秒 (Haiku)", sz=15, c=GREEN, bold=True)
    _p(tf, "节省 50%+ 延迟 · 用户感知连贯", sz=13, c=SUB)

    # 4 optimization techniques
    _txt(s, 0.7, 3.8, 8, 0.4, "4 项关键优化", sz=18, c=TITLE, bold=True)

    opts = [
        ("并行记忆预取", "ASR 完成后立即启动向量检索\n与意图路由同时执行", "节省 ~3s", GREEN),
        ("流式 TTS", "LLM 出一句话就立即合成\n不等全部生成完", "首句提前 5s+", BLUE),
        ("模型降级链", "Opus → Sonnet → Haiku\n自动切换保证可用性", "TTFT 5s→1s", PURPLE),
        ("本地语音引擎", "端侧 ASR+TTS 推理\n无网络延迟", "0.5s vs 3s", ORANGE),
    ]
    for i, (name, desc, effect, clr) in enumerate(opts):
        x = 0.5 + i * 3.15
        bg = {GREEN: GREEN_BG, BLUE: BLUE_BG, PURPLE: PURPLE_BG, ORANGE: ORANGE_BG}[clr]
        _card(s, x, 4.3, 2.95, 2.5, fill=bg)
        _txt(s, x + 0.15, 4.4, 2.65, 0.3, name, sz=14, c=clr, bold=True)
        lines = desc.split("\n")
        _txt(s, x + 0.15, 4.8, 2.65, 0.3, lines[0], sz=12, c=BODY)
        if len(lines) > 1:
            _txt(s, x + 0.15, 5.1, 2.65, 0.3, lines[1], sz=12, c=SUB)
        _card(s, x + 0.15, 5.6, 2.65, 0.5, fill=clr)
        _txt(s, x + 0.15, 5.65, 2.65, 0.35, effect, sz=14, c=WHITE, bold=True, al=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 10 — 延迟瀑布图对比 (DATA)
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "数据: 延迟瀑布对比 — 传统 vs Askme",
            "基于实测数据 (Opus via relay, sunrise 设备)")

    # Waterfall comparison table
    stages_data = [
        ("语音识别 ASR", "3.0s (云端API)", "0.3-0.8s (端侧推理)", GREEN),
        ("记忆检索", "3.0s (串行等待)", "0s (与 ASR 并行)", GREEN),
        ("意图路由", "0.5s (LLM 判断)", "<10ms (规则匹配)", GREEN),
        ("LLM 首字节", "5.0s (Opus)", "1.0s (Haiku voice_model)", BLUE),
        ("TTS 合成", "2.0s (全文缓冲)", "0.5s (首句流式)", BLUE),
        ("用户等待总计", "~13.5 秒", "~2-3 秒 (Haiku)", RED),
    ]

    # Table header
    _card(s, 0.5, 1.3, 12.3, 0.5, fill=CARD)
    _txt(s, 0.7, 1.35, 3.0, 0.35, "处理阶段", sz=14, c=TITLE, bold=True)
    _txt(s, 4.0, 1.35, 3.5, 0.35, "传统串行方案", sz=14, c=RED, bold=True)
    _txt(s, 7.8, 1.35, 4.5, 0.35, "Askme 并行+流式", sz=14, c=GREEN, bold=True)

    y = 1.9
    for stage, trad, askme_v, clr in stages_data:
        is_last = stage.startswith("用户")
        bg_r = RED_BG if is_last else (WHITE if stages_data.index((stage, trad, askme_v, clr)) % 2 == 0 else CARD)
        _card(s, 0.5, y, 12.3, 0.5, fill=bg_r)
        _txt(s, 0.7, y + 0.08, 3.0, 0.3, stage, sz=13, c=TITLE, bold=is_last)
        _txt(s, 4.0, y + 0.08, 3.5, 0.3, trad, sz=13, c=RED if is_last else SUB)
        _txt(s, 7.8, y + 0.08, 4.5, 0.3, askme_v, sz=13, c=clr, bold=True)
        y += 0.52

    # Key insight callout
    _card(s, 0.5, 5.2, 6.0, 1.8, fill=GREEN_BG)
    _txt(s, 0.8, 5.3, 5.4, 0.35, "关键洞察: voice_model 双模型策略", sz=16, c=GREEN, bold=True)
    tf = _txt(s, 0.8, 5.75, 5.4, 0.3, "", sz=13, c=BODY)
    _p(tf, "思考模型: Opus 4.6 (深度推理，5s TTFT)", sz=13, c=BODY)
    _p(tf, "语音模型: Haiku 4.5 (快速响应，1s TTFT)", sz=13, c=BODY)
    _p(tf, "", sz=4)
    _p(tf, "管理后台一键切换语音模型:", sz=13, c=BODY)
    _p(tf, "  思考模式 ↔ 实时语音模式", sz=12, c=GREEN, bold=True)

    _card(s, 6.8, 5.2, 6.0, 1.8, fill=BLUE_BG)
    _txt(s, 7.1, 5.3, 5.4, 0.35, "延迟优化公式", sz=16, c=BLUE, bold=True)
    tf = _txt(s, 7.1, 5.75, 5.4, 0.3, "", sz=13, c=BODY)
    _p(tf, "传统: ASR + 记忆 + 路由 + LLM + TTS (串行)", sz=12, c=SUB)
    _p(tf, "= 3 + 3 + 0.5 + 5 + 2 = 13.5s", sz=13, c=RED, bold=True)
    _p(tf, "", sz=4)
    _p(tf, "Askme: max(ASR, 记忆) + 路由 + 流式(LLM+TTS)", sz=12, c=SUB)
    _p(tf, "= 0.8 + 0 + 0.01 + 1.0 + 0.5 = 2.3s", sz=13, c=GREEN, bold=True)

    # ══════════════════════════════════════════════════════════
    # SLIDE 11 — 记忆增长数据
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "数据: 记忆增长曲线 — 机器人怎么变聪明",
            "基于 8h/天巡检场景推算")

    # Growth table
    _card(s, 0.5, 1.3, 12.3, 0.5, fill=CARD)
    cols = ["指标", "Day 0", "Day 1", "Day 3", "Day 7", "Day 30"]
    cx = [0.7, 3.5, 5.2, 6.9, 8.6, 10.5]
    for x, h in zip(cx, cols):
        _txt(s, x, 1.35, 1.6, 0.35, h, sz=13, c=TITLE, bold=True)

    growth = [
        ("L1 对话轮次", ["0", "40 (滑动)", "40 (滑动)", "40 (滑动)", "40 (滑动)"], BLUE),
        ("L2 会话摘要", ["0", "3 条", "12 条", "28 条", "120 条"], PURPLE),
        ("L3 事件日志", ["0", "~200", "~600", "~1,400", "~6,000"], GREEN),
        ("L3 反思次数", ["0", "2-3 次", "8-10 次", "20+ 次", "90+ 次"], GREEN),
        ("L3 世界知识", ["0", "5-10 条", "25-40 条", "60-100 条", "200+ 条"], GREEN),
        ("知识类别覆盖", ["0/5", "2/5", "4/5", "5/5", "5/5"], CYAN),
    ]

    y = 1.9
    for metric, vals, clr in growth:
        bg_r = WHITE if growth.index((metric, vals, clr)) % 2 == 0 else CARD
        _card(s, 0.5, y, 12.3, 0.48, fill=bg_r)
        _txt(s, 0.7, y + 0.06, 2.5, 0.3, metric, sz=12, c=BODY)
        for i, v in enumerate(vals):
            is_last = i >= 3
            _txt(s, cx[i + 1], y + 0.06, 1.6, 0.3, v, sz=12,
                 c=clr if is_last else SUB, bold=is_last)
        y += 0.5

    # Calculation basis
    _card(s, 0.5, 4.9, 6.0, 2.3, fill=GREEN_BG)
    _txt(s, 0.8, 5.0, 5.4, 0.35, "推算依据", sz=16, c=GREEN, bold=True)
    tf = _txt(s, 0.8, 5.45, 5.4, 0.3, "", sz=12, c=BODY)
    _p(tf, "每次巡检指令: importance ~0.7", sz=12, c=BODY)
    _p(tf, "反思触发: 累积 importance ≥ 15 (~22 条指令)", sz=12, c=BODY)
    _p(tf, "反思冷却: 5 分钟，每天 ~3 次反思", sz=12, c=BODY)
    _p(tf, "每次反思产出: 3-8 条知识事实", sz=12, c=BODY)
    _p(tf, "知识类别: 环境/实体/规律/交互/自省", sz=12, c=BODY)
    _p(tf, "遗忘曲线: S_init=1h, 访问×2, 上限 7 天", sz=12, c=SUB)

    # What this means
    _card(s, 6.8, 4.9, 6.0, 2.3, fill=PURPLE_BG)
    _txt(s, 7.1, 5.0, 5.4, 0.35, "这意味着什么", sz=16, c=PURPLE, bold=True)
    tf = _txt(s, 7.1, 5.45, 5.4, 0.3, "", sz=13, c=BODY)
    _p(tf, "Day 1: 机器人是新手，什么都要问", sz=13, c=SUB)
    _p(tf, "Day 7: 认识了环境、人员、规律", sz=13, c=BODY)
    _p(tf, "Day 30: 成为这个场地的\"老员工\"", sz=13, c=PURPLE, bold=True)
    _p(tf, "", sz=6)
    _p(tf, "不需要重新训练 · 不需要手动标注", sz=13, c=BODY)
    _p(tf, "纯靠日常运行自动积累经验", sz=13, c=GREEN, bold=True)

    # ══════════════════════════════════════════════════════════
    # SLIDE 12 — 安全响应数据
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "数据: 安全响应时间 — 工业级要求",
            "ESTOP 路径 vs LLM 路径延迟对比")

    safety_data = [
        ("ESTOP 关键词检测", "无此概念", "<10ms (规则匹配)", RED),
        ("ESTOP → 硬件停止", "无此概念", "<100ms (bypass LLM)", RED),
        ("误听指令防护", "直接执行", "Intent Router 优先级拦截", ORANGE),
        ("危险操作确认", "无确认机制", "语音确认 (30s 超时自动拒绝)", ORANGE),
        ("离线急停", "不可能 (云端)", "可用 (本地 ASR + 规则)", BLUE),
        ("操作审计", "无", "telemetry-hub 全链路记录", CYAN),
        ("安全策略远程更新", "无", "OTA 推送 policy 更新", GREEN),
    ]

    _card(s, 0.5, 1.3, 12.3, 0.5, fill=CARD)
    _txt(s, 0.7, 1.35, 3.5, 0.35, "安全场景", sz=14, c=TITLE, bold=True)
    _txt(s, 4.5, 1.35, 3.0, 0.35, "云端语音助手", sz=14, c=RED, bold=True)
    _txt(s, 7.8, 1.35, 5.0, 0.35, "Askme", sz=14, c=GREEN, bold=True)

    y = 1.9
    for scenario, cloud, askme_v, clr in safety_data:
        bg_r = WHITE if safety_data.index((scenario, cloud, askme_v, clr)) % 2 == 0 else CARD
        _card(s, 0.5, y, 12.3, 0.48, fill=bg_r)
        _txt(s, 0.7, y + 0.06, 3.5, 0.3, scenario, sz=12, c=BODY)
        _txt(s, 4.5, y + 0.06, 3.0, 0.3, cloud, sz=12, c=RED)
        _txt(s, 7.8, y + 0.06, 5.0, 0.3, askme_v, sz=12, c=clr, bold=True)
        y += 0.5

    # ESTOP path diagram
    _card(s, 0.5, 5.5, 6.0, 1.6, fill=RED_BG)
    _txt(s, 0.8, 5.6, 5.4, 0.35, "ESTOP 路径 — 零 LLM 延迟", sz=15, c=RED, bold=True)
    tf = _txt(s, 0.8, 6.0, 5.4, 0.3, "", sz=12, c=BODY)
    _p(tf, "\"停！\" → ASR → IntentRouter (Tier 1) → dog-safety", sz=12, c=BODY)
    _p(tf, "跳过: LLM · 记忆检索 · Skill 匹配 · 工具调用", sz=12, c=RED)
    _p(tf, "总延迟: ASR 端点检测 + <10ms = <1s", sz=13, c=RED, bold=True)

    _card(s, 6.8, 5.5, 6.0, 1.6, fill=BLUE_BG)
    _txt(s, 7.1, 5.6, 5.4, 0.35, "为什么这很重要", sz=15, c=BLUE, bold=True)
    tf = _txt(s, 7.1, 6.0, 5.4, 0.3, "", sz=12, c=BODY)
    _p(tf, "工业安全标准要求: 急停响应 < 500ms", sz=12, c=BODY)
    _p(tf, "云端助手最快也要 2-5s (网络往返)", sz=12, c=BODY)
    _p(tf, "Askme 的架构保证 ESTOP 不经过任何云服务", sz=13, c=BLUE, bold=True)

    # ══════════════════════════════════════════════════════════
    # SLIDE 13 — 竞品对比 (真实行业竞品)
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "竞品对比 — 真实行业方案 vs Askme",
            "没有直接竞品 = 新品类机会")

    # 5 columns: 能力 | Picovoice | 讯飞/Riva | 厂商自研 | Askme
    headers = ["能力维度", "Picovoice", "讯飞/Riva", "厂商自研", "Askme"]
    hdr_x = [0.5, 2.7, 4.8, 7.0, 9.5]
    hdr_w = [2.0, 1.9, 2.0, 2.3, 3.5]

    _card(s, 0.3, 1.2, 12.7, 0.48, fill=CARD)
    for hx, hw, h in zip(hdr_x, hdr_w, headers):
        _txt(s, hx, 1.23, hw, 0.35, h, sz=13, c=TITLE, bold=True)

    rows = [
        ("定位",      "语音组件库",    "云端 API",    "固定指令集",   "机器人 AI Runtime", BLUE),
        ("语音 ESTOP", "需自行开发",   "无此概念",    "硬件按钮",     "语音<100ms 直达",  RED),
        ("安全审批",   "无",           "无",          "无",           "3 级+语音确认",    ORANGE),
        ("记忆/学习",  "无",           "无",          "无",           "4层+反思+遗忘曲线", PURPLE),
        ("离线能力",   "全离线 (优势)", "完全依赖云",  "全离线",       "全离线+云端增强",  GREEN),
        ("中文优化",   "一般",         "优秀 (讯飞)", "不适用",       "中文深度优化",     BLUE),
        ("技能扩展",   "需编程",       "需编程",      "固件重写",     "零代码技能包",     GREEN),
        ("部署成本",   "按设备$年费",  "按调用量",    "自研高投入",   "Docker 本地",     CYAN),
        ("适用场景",   "消费电子",     "客服/车载",   "特定产品",     "所有工业机器人",   PURPLE),
    ]

    y = 1.75
    for i, (label, pico, xf, oem, askme_v, clr) in enumerate(rows):
        bg_r = WHITE if i % 2 == 0 else CARD
        _card(s, 0.3, y, 12.7, 0.44, fill=bg_r)
        _txt(s, 0.5, y + 0.05, 2.0, 0.3, label, sz=11, c=BODY)
        _txt(s, 2.7, y + 0.05, 1.9, 0.3, pico, sz=11, c=MUTED)
        _txt(s, 4.8, y + 0.05, 2.0, 0.3, xf, sz=11, c=MUTED)
        _txt(s, 7.0, y + 0.05, 2.3, 0.3, oem, sz=11, c=MUTED)
        _txt(s, 9.5, y + 0.05, 3.5, 0.3, askme_v, sz=11, c=clr, bold=True)
        y += 0.46

    # Bottom insight
    _card(s, 0.3, 6.0, 12.7, 0.6, fill=BLUE_BG)
    _txt(s, 0.5, 6.05, 12.3, 0.45,
         "Picovoice 做语音组件 · 讯飞/Riva 做云端 API · 厂商自研只解决自己 → 没有人做\"机器人语音 AI 运行时\"",
         sz=13, c=BLUE, bold=True, al=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 14 — GTM 推广策略
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "GTM 策略 — 怎么把 Askme 卖出去",
            "早期聚焦 · 做深一个场景 · 用案例复制")

    # 4 strategy cards
    gtm = [
        ("P0 目标客户", BLUE, BLUE_BG, [
            "电力/化工巡检集成商",
            "  亿嘉和·国自·朗驰欣创",
            "安全刚需+离线刚需+客单高",
            "  机器人 15-40 万，语音<5%",
            "LingTu 存量客户反向导入",
            "  导航+语音捆绑推",
        ]),
        ("定价 Open Core", GREEN, GREEN_BG, [
            "社区版: 免费开源核心",
            "  吸引开发者，建立生态",
            "专业版: ¥8k-15k/设备/年",
            "  全功能+6微服务 runtime",
            "企业版: 项目制 ¥5-10万",
            "  定制部署+技能包+支持",
        ]),
        ("切入策略", PURPLE, PURPLE_BG, [
            "不等 Thunder，先跑模拟器",
            "  Gazebo+LingTu+Askme Demo",
            "录制 60s Demo 视频",
            "  比 50 页 PPT 有说服力",
            "标杆客户: 半价换案例权",
            "  用案例撬动行业",
        ]),
        ("竞争壁垒", ORANGE, ORANGE_BG, [
            "安全系统: 讯飞/Picovoice 不做",
            "  ESTOP+三级审批=真差异",
            "全栈: Askme+LingTu+Brainstem",
            "  从语音到动作完整闭环",
            "数据: 越用越聪明",
            "  客户换不掉积累的经验",
        ]),
    ]
    x = 0.3
    for title, clr, bg, items in gtm:
        _card(s, x, 1.2, 3.1, 3.8, fill=bg)
        _txt(s, x + 0.15, 1.3, 2.8, 0.3, title, sz=15, c=clr, bold=True)
        ty = 1.7
        for item in items:
            is_indent = item.startswith("  ")
            _txt(s, x + (0.3 if is_indent else 0.15), ty, 2.65, 0.25,
                 item.strip(), sz=11 if is_indent else 12,
                 c=SUB if is_indent else BODY,
                 bold=not is_indent)
            ty += 0.28
        x += 3.25

    # Timeline
    _card(s, 0.3, 5.25, 12.7, 2.0, fill=CARD)
    _txt(s, 0.5, 5.35, 12.3, 0.35, "12 个月里程碑", sz=16, c=TITLE, bold=True)

    timeline = [
        ("M1-3", "Demo + 标杆", GREEN, GREEN_BG,
         "录制 Demo 视频\n找到 1 个试用客户\n完善文档+快速部署脚本"),
        ("M4-6", "签单 + 复制", BLUE, BLUE_BG,
         "签第 1 个付费合同\nLingTu+Askme 整机方案\n完善技能包生态"),
        ("M7-9", "规模化", PURPLE, PURPLE_BG,
         "3-5 个付费客户\n开放技能包市场\nFleet 记忆共享功能"),
        ("M10-12", "生态建设", ORANGE, ORANGE_BG,
         "开发者社区\n集成商渠道合作\n海外市场验证"),
    ]
    x = 0.4
    for phase, title, clr, bg, desc in timeline:
        _card(s, x, 5.8, 3.0, 1.25, fill=bg)
        _txt(s, x + 0.1, 5.85, 0.6, 0.25, phase, sz=12, c=clr, bold=True)
        _txt(s, x + 0.7, 5.85, 2.2, 0.25, title, sz=12, c=TITLE, bold=True)
        lines = desc.split("\n")
        ty = 6.15
        for line in lines:
            _txt(s, x + 0.1, ty, 2.8, 0.22, line, sz=10, c=BODY)
            ty += 0.22
        if x < 9:
            _txt(s, x + 2.9, 6.1, 0.3, 0.3, "→", sz=14, c=clr, bold=True)
        x += 3.15

    # ══════════════════════════════════════════════════════════
    # SLIDE 15 — 下一步
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "下一步 — 从 Demo 到真机")

    milestones = [
        ("立即可做", GREEN, GREEN_BG, [
            "nav-gateway 接入 LingTu gRPC (v1.7.5 稳定)",
            "部署端侧语音模型到边缘设备",
            "语音控制 Thunder: \"去厨房\" → 导航执行",
        ]),
        ("等硬件解锁", ORANGE, ORANGE_BG, [
            "han_dog.dart:191 电机输出取消注释",
            "dog-control 接入 Brainstem gRPC",
            "语音控制: \"坐下\" \"走\" → 真实运动",
        ]),
        ("商业化方向", BLUE, BLUE_BG, [
            "LingTu + Askme 打包方案 (最快签单路径)",
            "Skill 市场 — 不同客户不同技能包",
            "Fleet 记忆共享 — 多机器人共享学习经验",
        ]),
    ]
    x = 0.5
    for title, clr, bg, items in milestones:
        _card(s, x, 1.2, 4.0, 3.5, fill=bg)
        _txt(s, x + 0.2, 1.3, 3.6, 0.35, title, sz=18, c=clr, bold=True)
        y = 1.8
        for item in items:
            _txt(s, x + 0.2, y, 3.6, 0.55, f"• {item}", sz=13, c=BODY)
            y += 0.7
        x += 4.15

    # Value proposition
    _card(s, 0.5, 5.0, 12.3, 2.0, fill=BLUE_BG)
    _txt(s, 0.8, 5.15, 11.5, 0.5,
         "Askme 的核心价值", sz=20, c=BLUE, bold=True)
    _txt(s, 0.8, 5.7, 11.5, 0.8,
         "一套面向工业机器人的语音 AI 运行时——\n"
         "让 LLM 变成可信赖的现场 Agent: 听得懂、记得住、管得严、接得快",
         sz=18, c=TITLE)
    _txt(s, 0.8, 6.5, 11.5, 0.4,
         "穹沛科技 · 让机器人听懂你", sz=14, c=MUTED, al=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────
    prs.save(out)
    print(f"PPT saved: {out}")


if __name__ == "__main__":
    build("docs/NOVA_Dog_Askme_v4_Report_v4.pptx")
