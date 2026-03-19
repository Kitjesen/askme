"""Build a demo PPT with embedded voice audio showing the complete inspection scenario."""

import os
import shutil
import zipfile
import tempfile
from copy import deepcopy
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI

# ── Palette ──────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG         = RGBColor(0xFA, 0xFA, 0xFC)
CARD       = RGBColor(0xF0, 0xF2, 0xF5)
TITLE      = RGBColor(0x11, 0x18, 0x27)
BODY       = RGBColor(0x33, 0x40, 0x55)
SUB        = RGBColor(0x64, 0x74, 0x8B)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
BLUE       = RGBColor(0x25, 0x63, 0xEB)
BLUE_BG    = RGBColor(0xEF, 0xF6, 0xFF)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_BG  = RGBColor(0xF5, 0xF3, 0xFF)
GREEN      = RGBColor(0x05, 0x96, 0x69)
GREEN_BG   = RGBColor(0xEC, 0xFD, 0xF5)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_BG  = RGBColor(0xFF, 0xF7, 0xED)
RED        = RGBColor(0xDC, 0x26, 0x26)
RED_BG     = RGBColor(0xFE, 0xF2, 0xF2)
CYAN       = RGBColor(0x06, 0x84, 0x8D)
CYAN_BG    = RGBColor(0xF0, 0xFD, 0xFA)

FONT = "Microsoft YaHei"
W = 13.333
H = 7.5


def _bg(slide, color=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _txt(slide, l, t, w, h, text, sz=16, c=BODY, bold=False, al=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = FONT; p.alignment = al
    return tf

def _p(tf, text, sz=14, c=BODY, bold=False, sp=Pt(4)):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = FONT; p.space_before = sp
    return p

def _card(slide, l, t, w, h, fill=CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s

def _bar(slide, l, t, w, c=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def _header(slide, text, sub=None):
    _bar(slide, 0.7, 0.4, 0.5, BLUE)
    _txt(slide, 1.35, 0.28, 10, 0.6, text, sz=28, c=TITLE, bold=True)
    if sub:
        _txt(slide, 1.35, 0.78, 10, 0.4, sub, sz=15, c=SUB)

def _play_btn(slide, l, t, w, h, clr, label=""):
    """Draw a play button indicator (triangle in circle)."""
    # Circle background
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    circ.fill.solid(); circ.fill.fore_color.rgb = clr
    circ.line.fill.background()
    # Play triangle text
    tf = circ.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "▶"
    p.font.size = Pt(int(h * 20))
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return circ


def build_demo_ppt(output_path, speaker_dir="docs/voice_samples/speaker_10"):
    """Build a standalone demo PPT with voice samples."""
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    B = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════
    # SLIDE 1 — 产品方案总览
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _bar(s, 0, 0, W, BLUE)
    _txt(s, 1.0, 0.6, 10, 0.8,
         "Askme 巡检导航方案", sz=40, c=TITLE, bold=True)
    _txt(s, 1.0, 1.3, 10, 0.5,
         "语音操控 + 自主导航 + 安全审批 + 持续学习 = 开箱即用的巡检机器人大脑",
         sz=18, c=BLUE, bold=True)

    # 4 product modules
    modules = [
        ("Askme 语音 AI", "语音交互运行时",
         "语音唤醒 → 识别 → 意图路由 → LLM 对话\n"
         "流式 TTS 即时播报 · 三级安全审批\n"
         "四层记忆持续学习 · 12+ 内置技能",
         BLUE, BLUE_BG),
        ("LingTu 导航", "自主导航引擎 v1.7.5",
         "SLAM 建图 · 路径规划 · 避障\n"
         "6 种导航模式 · 多楼层支持\n"
         "gRPC 接口 · 已稳定运行",
         GREEN, GREEN_BG),
        ("NOVA Runtime", "6 微服务编排",
         "arbiter 任务状态机\n"
         "dog-safety 安全守卫 · telemetry 遥测\n"
         "Docker Compose 一键部署",
         PURPLE, PURPLE_BG),
        ("OTA 远程运维", "空中升级基础设施",
         "远程推送新技能 · 新安全策略\n"
         "设备健康监控 · 固件更新\n"
         "不用派人到现场",
         ORANGE, ORANGE_BG),
    ]
    x = 0.4
    for title, sub, desc, clr, bg in modules:
        _card(s, x, 2.0, 3.05, 2.8, fill=bg)
        _txt(s, x + 0.15, 2.1, 2.75, 0.3, title, sz=16, c=clr, bold=True)
        _txt(s, x + 0.15, 2.45, 2.75, 0.25, sub, sz=12, c=SUB)
        _bar(s, x + 0.15, 2.75, 1.2, clr)
        lines = desc.split("\n")
        ty = 2.9
        for line in lines:
            _txt(s, x + 0.15, ty, 2.75, 0.25, line, sz=11, c=BODY)
            ty += 0.28
        x += 3.2

    # Use case scenarios
    _bar(s, 0.7, 5.1, 0.4, BLUE)
    _txt(s, 1.2, 5.0, 8, 0.4, "适用场景", sz=18, c=TITLE, bold=True)

    scenes = [
        ("电力巡检", "变电站/配电房/线路", "语音下令巡检路线\n异常自动报告", BLUE),
        ("化工安防", "管廊/罐区/危险区域", "离线运行\n语音急停保障安全", RED),
        ("仓储物流", "货架/通道/月台", "语音查库存\n任务语音派发", GREEN),
        ("园区安保", "楼宇/停车场/周界", "定时巡逻\n发现异常语音报警", PURPLE),
    ]
    x = 0.4
    for title, area, desc, clr in scenes:
        bg = {BLUE: BLUE_BG, RED: RED_BG, GREEN: GREEN_BG, PURPLE: PURPLE_BG}[clr]
        _card(s, x, 5.5, 3.05, 1.6, fill=bg)
        _txt(s, x + 0.15, 5.55, 1.2, 0.25, title, sz=14, c=clr, bold=True)
        _txt(s, x + 1.3, 5.55, 1.6, 0.25, area, sz=10, c=SUB)
        lines = desc.split("\n")
        _txt(s, x + 0.15, 5.9, 2.75, 0.25, lines[0], sz=11, c=BODY)
        if len(lines) > 1:
            _txt(s, x + 0.15, 6.15, 2.75, 0.25, lines[1], sz=11, c=BODY)
        x += 3.2

    # ══════════════════════════════════════════════════════════
    # SLIDE 2 — 巡检演示: 语音效果 (上半部)
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "实际语音效果 — 点击 ▶ 播放真实 TTS 音频",
            "VITS-Aishell3 本地离线合成 · RTF 0.2 (5x 实时) · 零网络依赖")

    # Demo steps with audio indicators
    demo_steps = [
        ("① 唤醒+问候", "操作员:", "\"小穹\"  (唤醒词触发)",
         "Thunder:", "你好，我是Thunder，穹沛的巡检机器人。等待指令。",
         "01_greeting.wav", BLUE, BLUE_BG),

        ("② 下达任务", "操作员:", "\"去三号配电柜巡检\"",
         "Thunder:", "收到，正在前往三号配电柜巡检。预计到达时间两分钟。",
         "02_confirm_task.wav", GREEN, GREEN_BG),

        ("③ 安全审批", "操作员:", "\"把机械臂移到前面\"",
         "Thunder:", "这是高风险操作：移动机械臂到前方位置。请说确认执行继续，或取消放弃。",
         "04_safety_confirm.wav", ORANGE, ORANGE_BG),
    ]

    y = 1.2
    for step_title, speaker1, text1, speaker2, text2, wav, clr, bg in demo_steps:
        _card(s, 0.3, y, 12.7, 1.7, fill=bg)

        # Step title
        _txt(s, 0.5, y + 0.05, 2.5, 0.3, step_title, sz=15, c=clr, bold=True)

        # Operator line
        _txt(s, 0.5, y + 0.4, 1.0, 0.25, speaker1, sz=12, c=BLUE, bold=True)
        _txt(s, 1.5, y + 0.4, 5.5, 0.25, text1, sz=12, c=BODY)

        # Robot response line
        _txt(s, 0.5, y + 0.75, 1.0, 0.25, speaker2, sz=12, c=clr, bold=True)
        _txt(s, 1.5, y + 0.75, 8.5, 0.5, text2, sz=12, c=TITLE)

        # Play button + audio file label
        _play_btn(s, 10.8, y + 0.35, 0.7, 0.7, clr)
        _txt(s, 10.5, y + 1.15, 2.5, 0.25, f"▲ {wav}", sz=9, c=SUB, al=PP_ALIGN.CENTER)

        # System component tag
        tags = {"01": "KWS → ASR", "02": "Intent → Skill", "04": "Safety 审批"}
        tag_key = wav[:2]
        if tag_key in tags:
            _txt(s, 7.5, y + 0.05, 5.0, 0.25, tags[tag_key], sz=11, c=clr, al=PP_ALIGN.RIGHT)

        y += 1.85

    # ══════════════════════════════════════════════════════════
    # SLIDE 3 — 巡检演示: 语音效果 (下半部)
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "实际语音效果 (续) — 异常处理 · 急停 · 记忆",
            "展示安全响应和持续学习能力")

    demo_steps_2 = [
        ("④ 巡检报告", "Thunder:", "(到达目标，自动检测)",
         "Thunder:", "巡逻报告：三号区域检测到温度异常，柜体表面六十八度，超过阈值。已记录并上报工单。",
         "05_patrol_report.wav", PURPLE, PURPLE_BG),

        ("⑤ 紧急停止", "操作员:", "\"停！\"  (ESTOP 触发，<100ms)",
         "Thunder:", "已紧急停机。所有运动已暂停，等待下一步指令。",
         "03_estop.wav", RED, RED_BG),

        ("⑥ 记忆调用", "操作员:", "\"之前这里有什么异常吗？\"",
         "Thunder:", "根据之前的巡检记录，这个区域上周也出现过类似温度异常。建议排查散热系统。",
         "06_memory_demo.wav", CYAN, CYAN_BG),
    ]

    y = 1.2
    for step_title, speaker1, text1, speaker2, text2, wav, clr, bg in demo_steps_2:
        _card(s, 0.3, y, 12.7, 1.7, fill=bg)
        _txt(s, 0.5, y + 0.05, 2.5, 0.3, step_title, sz=15, c=clr, bold=True)
        _txt(s, 0.5, y + 0.4, 1.0, 0.25, speaker1, sz=12, c=BLUE if "操作" in speaker1 else clr, bold=True)
        _txt(s, 1.5, y + 0.4, 5.5, 0.25, text1, sz=12, c=BODY)
        _txt(s, 0.5, y + 0.75, 1.0, 0.25, speaker2, sz=12, c=clr, bold=True)
        _txt(s, 1.5, y + 0.75, 8.5, 0.5, text2, sz=12, c=TITLE)
        _play_btn(s, 10.8, y + 0.35, 0.7, 0.7, clr)
        _txt(s, 10.5, y + 1.15, 2.5, 0.25, f"▲ {wav}", sz=9, c=SUB, al=PP_ALIGN.CENTER)

        tags = {"05": "Skill → Telemetry", "03": "ESTOP Tier-1", "06": "L3 情景记忆"}
        tag_key = wav[:2]
        if tag_key in tags:
            _txt(s, 7.5, y + 0.05, 5.0, 0.25, tags[tag_key], sz=11, c=clr, al=PP_ALIGN.RIGHT)

        y += 1.85

    # Bottom: TTS performance data
    _card(s, 0.3, 6.65, 12.7, 0.55, fill=CARD)
    _txt(s, 0.5, 6.7, 12.3, 0.4,
         "TTS 性能: 本地 VITS 合成 · RTF=0.2 (1秒音频只需 0.2秒生成) · "
         "短句 0.9s / 长句 1.8s · 纯 CPU · 零网络 · 329MB 模型",
         sz=12, c=SUB, al=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 4 — 完整方案架构
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(B); _bg(s)
    _header(s, "巡检导航完整方案 — 从部署到运行",
            "Askme + LingTu + NOVA Runtime 一站式交付")

    # Deployment flow
    deploy = [
        ("第1步: 部署", "Docker Compose\n一键启动 6 个微服务", GREEN, GREEN_BG),
        ("第2步: 配置", "管理后台设定\n人格·技能·安全策略", BLUE, BLUE_BG),
        ("第3步: 建图", "LingTu SLAM\n激光雷达自动建图", PURPLE, PURPLE_BG),
        ("第4步: 运行", "语音+导航+巡检\n开箱即用", ORANGE, ORANGE_BG),
    ]
    x = 0.3
    for title, desc, clr, bg in deploy:
        _card(s, x, 1.2, 3.05, 1.4, fill=bg)
        _txt(s, x + 0.15, 1.3, 2.75, 0.3, title, sz=14, c=clr, bold=True)
        lines = desc.split("\n")
        _txt(s, x + 0.15, 1.65, 2.75, 0.25, lines[0], sz=12, c=BODY)
        if len(lines) > 1:
            _txt(s, x + 0.15, 1.9, 2.75, 0.25, lines[1], sz=11, c=SUB)
        if x < 9:
            _txt(s, x + 2.95, 1.6, 0.3, 0.3, "→", sz=16, c=clr, bold=True)
        x += 3.2

    # What operator can do
    _bar(s, 0.7, 2.9, 0.4, BLUE)
    _txt(s, 1.2, 2.8, 8, 0.4, "操作员能做什么 — 全部通过语音", sz=18, c=TITLE, bold=True)

    capabilities = [
        ("巡检派令", "\"去三号配电柜\"", "nav-gateway\n→ LingTu 执行", GREEN),
        ("状态查询", "\"电池还有多少\"", "telemetry-hub\n→ 实时数据", BLUE),
        ("异常处理", "\"拍照记录\"", "skill: 环境感知\n→ 自动上报", PURPLE),
        ("紧急制动", "\"停！\"", "dog-safety\nESTOP <100ms", RED),
        ("任务管理", "\"取消当前任务\"", "arbiter\n→ 状态机", ORANGE),
        ("经验查询", "\"上次这里有问题吗\"", "L3 情景记忆\n→ 知识检索", CYAN),
    ]
    for i, (cap, example, system, clr) in enumerate(capabilities):
        col = i % 3
        row = i // 3
        x = 0.3 + col * 4.3
        y = 3.3 + row * 1.75
        bg = {GREEN: GREEN_BG, BLUE: BLUE_BG, PURPLE: PURPLE_BG,
              RED: RED_BG, ORANGE: ORANGE_BG, CYAN: CYAN_BG}[clr]
        _card(s, x, y, 4.0, 1.5, fill=bg)
        _txt(s, x + 0.15, y + 0.1, 1.5, 0.3, cap, sz=14, c=clr, bold=True)
        _txt(s, x + 1.7, y + 0.1, 2.2, 0.3, example, sz=12, c=BODY)
        lines = system.split("\n")
        _txt(s, x + 0.15, y + 0.5, 3.7, 0.25, lines[0], sz=11, c=SUB)
        if len(lines) > 1:
            _txt(s, x + 0.15, y + 0.75, 3.7, 0.25, lines[1], sz=11, c=SUB)

        # Processing time indicator
        times = {"紧急制动": "<100ms", "状态查询": "~1s", "巡检派令": "~2s",
                 "异常处理": "~3s", "任务管理": "~1s", "经验查询": "~2s"}
        _txt(s, x + 0.15, y + 1.1, 3.7, 0.25,
             f"响应: {times.get(cap, '')}", sz=10, c=clr, bold=True)

    # ── Save base PPTX ──────────────────────────────────────
    prs.save(output_path)
    print(f"Base demo PPT saved: {output_path}")
    return output_path


def embed_audio_files(pptx_path, speaker_dir, output_path):
    """Post-process PPTX to embed WAV audio files into slides."""
    # Audio files to embed per slide (slide index → list of wav files)
    # Slide 2 (index 1): greeting, confirm_task, safety_confirm
    # Slide 3 (index 2): patrol_report, estop, memory_demo
    slide_audio = {
        1: [
            ("01_greeting.wav", 10.8, 0.35 + 1.2),
            ("02_confirm_task.wav", 10.8, 0.35 + 1.2 + 1.85),
            ("04_safety_confirm.wav", 10.8, 0.35 + 1.2 + 1.85 * 2),
        ],
        2: [
            ("05_patrol_report.wav", 10.8, 0.35 + 1.2),
            ("03_estop.wav", 10.8, 0.35 + 1.2 + 1.85),
            ("06_memory_demo.wav", 10.8, 0.35 + 1.2 + 1.85 * 2),
        ],
    }

    # Work on a temp copy
    tmp = pptx_path + ".tmp"
    shutil.copy2(pptx_path, tmp)

    try:
        with zipfile.ZipFile(tmp, 'r') as zin, \
             zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:

            # Copy all existing entries
            existing = set()
            for item in zin.infolist():
                data = zin.read(item.filename)

                # Patch [Content_Types].xml to add audio/wav
                if item.filename == '[Content_Types].xml':
                    ct_tree = etree.fromstring(data)
                    ns = 'http://schemas.openxmlformats.org/package/2006/content-types'
                    # Check if .wav extension already registered
                    has_wav = any(
                        el.get('Extension') == 'wav'
                        for el in ct_tree.findall(f'{{{ns}}}Default')
                    )
                    if not has_wav:
                        ext_el = etree.SubElement(ct_tree, f'{{{ns}}}Default')
                        ext_el.set('Extension', 'wav')
                        ext_el.set('ContentType', 'audio/wav')
                    data = etree.tostring(ct_tree, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)

                zout.writestr(item, data)
                existing.add(item.filename)

            # Add audio files to ppt/media/
            audio_idx = 1
            slide_rels_patches = {}  # slide_rels_path -> list of (rid, media_name)

            for slide_idx, audio_list in slide_audio.items():
                slide_num = slide_idx + 1
                rels_path = f'ppt/slides/_rels/slide{slide_num}.xml.rels'

                for wav_name, ax, ay in audio_list:
                    wav_path = os.path.join(speaker_dir, wav_name)
                    if not os.path.exists(wav_path):
                        print(f"  WARNING: {wav_path} not found, skipping")
                        continue

                    media_name = f'audio{audio_idx}.wav'
                    media_path = f'ppt/media/{media_name}'

                    # Add the WAV file
                    with open(wav_path, 'rb') as f:
                        zout.writestr(media_path, f.read())

                    if rels_path not in slide_rels_patches:
                        slide_rels_patches[rels_path] = []
                    slide_rels_patches[rels_path].append((media_name, audio_idx))

                    audio_idx += 1

            # Patch slide relationship files to add audio references
            # We need to re-open and patch. Since we already wrote, we need
            # to handle this differently. Let's use a second pass.
            print(f"  Embedded {audio_idx - 1} audio files into {output_path}")

    finally:
        if os.path.exists(tmp):
            os.remove(tmp)

    # Second pass: patch slide .rels files
    _patch_slide_rels(output_path, slide_rels_patches)

    print(f"Demo PPT with audio: {output_path}")


def _patch_slide_rels(pptx_path, rels_patches):
    """Patch slide .rels files to add audio relationships."""
    if not rels_patches:
        return

    tmp = pptx_path + ".patch"
    ns_rels = 'http://schemas.openxmlformats.org/package/2006/relationships'
    ns_audio = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio'
    ns_media = 'http://schemas.microsoft.com/office/2007/relationships/media'

    with zipfile.ZipFile(pptx_path, 'r') as zin, \
         zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:

        for item in zin.infolist():
            data = zin.read(item.filename)

            if item.filename in rels_patches:
                tree = etree.fromstring(data)

                # Find max existing rId
                max_rid = 0
                for rel in tree.findall(f'{{{ns_rels}}}Relationship'):
                    rid_str = rel.get('Id', 'rId0')
                    try:
                        num = int(rid_str.replace('rId', ''))
                        max_rid = max(max_rid, num)
                    except ValueError:
                        pass

                for media_name, audio_idx in rels_patches[item.filename]:
                    # Audio relationship
                    max_rid += 1
                    rel_audio = etree.SubElement(tree, f'{{{ns_rels}}}Relationship')
                    rel_audio.set('Id', f'rId{max_rid}')
                    rel_audio.set('Type', ns_audio)
                    rel_audio.set('Target', f'../media/{media_name}')

                    # Media relationship
                    max_rid += 1
                    rel_media = etree.SubElement(tree, f'{{{ns_rels}}}Relationship')
                    rel_media.set('Id', f'rId{max_rid}')
                    rel_media.set('Type', ns_media)
                    rel_media.set('Target', f'../media/{media_name}')

                data = etree.tostring(tree, xml_declaration=True,
                                      encoding='UTF-8', standalone=True)

            zout.writestr(item, data)

    # Replace original with patched
    os.replace(tmp, pptx_path)


if __name__ == "__main__":
    speaker_dir = "docs/voice_samples/speaker_10"
    base_path = "docs/Askme_Demo_Inspection.pptx"

    # Step 1: Build slides
    build_demo_ppt(base_path, speaker_dir)

    # Step 2: Embed audio
    embed_audio_files(base_path, speaker_dir, base_path)

    print("\nDone! Open in PowerPoint to test audio playback.")
    print("NOTE: Audio files are embedded in ppt/media/. In PowerPoint,")
    print("you may need to insert them as media objects for click-to-play.")
    print(f"\nAlternative: Play WAV files directly from {speaker_dir}/")
